from pptx import Presentation
from pptx.util import Inches

# Criar uma nova apresentação
prs = Presentation()

# Slide 1: Título e Componentes do Projeto
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
title.text = "Documentação do Programa de Cadastro de Alunos e Professores (OOP)"
subtitle = slide_title.placeholders[1]
subtitle.text = "Componentes: Ivan Varella e Ricardo Nogueira"

# Slide 2: Visão Geral
slide_visao_geral = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_visao_geral.shapes.title
title.text = "Visão Geral"
content = slide_visao_geral.placeholders[1].text_frame
content.text = (
    "O programa de Gerenciamento de Alunos e Professores é uma aplicação desenvolvida em Python que "
    "permite o cadastro, listagem, pesquisa e gerenciamento de informações de alunos e professores. "
    "O programa utiliza arquivos JSON para armazenar os dados, gera arquivos HTML e PDF com as informações do arquivo JSON."
)

# Slide 3: Estrutura do Código
slide_estrutura = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_estrutura.shapes.title
title.text = "Estrutura do Código"
content = slide_estrutura.placeholders[1].text_frame
content.text = (
    "O programa é composto pelos seguintes arquivos:\n\n"
    "├── aluno.py\n"
    "├── dados.html\n"
    "├── dados.json\n"
    "├── dados.pdf\n"
    "├── funcoesSuporte.py\n"
    "├── jsonHandler.py\n"
    "├── main.py\n"
    "├── pessoa.py\n"
    "├── professor.py\n"
    "└── rich_menus.py"
)

# Slide 4: Descrição dos Arquivos
slide_descr_files = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_descr_files.shapes.title
title.text = "Descrição dos Arquivos"
content = slide_descr_files.placeholders[1].text_frame
content.text = (
    "• main.py: Controla o fluxo do programa e exibe o menu principal.\n"
    "• funcoesSuporte.py: Contém funções de suporte e lógica do programa.\n"
    "• pessoa.py: Define a classe Pessoa, base para Aluno e Professor.\n"
    "• aluno.py: Define a classe Aluno, que herda da classe Pessoa.\n"
    "• professor.py: Define a classe Professor, que herda da classe Pessoa.\n"
    "• jsonHandler.py: Contém métodos para manipulação de dados em JSON.\n"
    "• dados.json: Armazena dados de alunos e professores em JSON.\n"
    "• rich_menus.py: Exibe o menu principal utilizando a biblioteca rich."
)

# Slide 5: Funcionamento do Programa
slide_funcionamento = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_funcionamento.shapes.title
title.text = "Funcionamento do Programa"
content = slide_funcionamento.placeholders[1].text_frame
content.text = (
    "• O programa é iniciado pelo arquivo main.py, que exibe o menu principal.\n"
    "• O usuário pode selecionar opções como cadastrar, listar, alterar e pesquisar alunos e professores.\n"
    "• Os dados são armazenados e lidos do arquivo dados.json usando a classe JsonHandler.\n"
    "• Relatórios HTML e PDF são gerados a partir dos dados em JSON."
)

# Slide 6: Fluxo da Lógica do Programa
slide_fluxo = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_fluxo.shapes.title
title.text = "Fluxo da Lógica do Programa"
content = slide_fluxo.placeholders[1].text_frame
content.text = (
    "• O código está organizado de forma modular, com cada arquivo responsável por uma parte específica.\n"
    "• Classes Pessoa, Aluno e Professor encapsulam informações e comportamentos.\n"
    "• A classe JsonHandler é responsável pelo CRUD no arquivo JSON.\n"
    "• A lógica do programa é centralizada no arquivo funcoesSuporte.py."
)

# Slide 7: Organização do Código
slide_organizacao = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_organizacao.shapes.title
title.text = "Organização do Código"
content = slide_organizacao.placeholders[1].text_frame
content.text = (
    "• Modularidade: Cada arquivo é responsável por uma parte do programa.\n"
    "• Encapsulamento: Classes encapsulam informações e comportamentos relacionados.\n"
    "• JsonHandler: Responsável pelo CRUD no arquivo JSON.\n"
    "• funcoesSuporte.py: Centraliza a lógica do programa e funções de suporte."
)

# Slide 8: Bibliotecas Utilizadas
slide_bibliotecas = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_bibliotecas.shapes.title
title.text = "Bibliotecas Utilizadas"
content = slide_bibliotecas.placeholders[1].text_frame
content.text = (
    "• Padrão do Python:\n"
    "  - os: Interface com o sistema operacional.\n"
    "  - copy: Criação de cópias profundas de objetos.\n"
    "  - platform: Informações sobre a plataforma.\n"
    "  - json: Manipulação de dados JSON.\n"
    "  - contextlib: Criação de gerenciadores de contexto personalizados.\n"
    "  - typing: Dicas de tipo para funções geradoras.\n"
    "  - time: Funções relacionadas ao tempo.\n"
    "• Não padrão do Python:\n"
    "  - reportlab: Geração de documentos PDF.\n"
    "  - rich: Criação de menus e interfaces de usuário formatadas."
)

# Salvar a apresentação no diretório atual
presentation_path = "Apresentacao_Cadastro_Alunos_Professores.pptx"
prs.save(presentation_path)

print(f"Apresentação salva em {presentation_path}")
